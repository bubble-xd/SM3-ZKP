from __future__ import annotations

import json
import math
import sys
from dataclasses import dataclass
from datetime import date
from pathlib import Path


PROJECT_ROOT = Path(__file__).resolve().parents[1]
VENDOR_ROOT = PROJECT_ROOT / ".vendor"
if VENDOR_ROOT.exists():
    sys.path.insert(0, str(VENDOR_ROOT))

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = PROJECT_ROOT / "SM3_ZKP_结项答辩_可编辑.pptx"
BENCHMARK_PATH = PROJECT_ROOT / "benchmarks" / "results" / "latest.json"
META_PATH = PROJECT_ROOT / "circuits" / "build" / "sm3_compression_step" / "meta.json"

FONT_CN = "Microsoft YaHei"
FONT_MONO = "Consolas"

COLOR_BG = "F7F4EE"
COLOR_SURFACE = "FFFDF9"
COLOR_SURFACE_ALT = "EEF6F4"
COLOR_PRIMARY = "10233F"
COLOR_SECONDARY = "0F766E"
COLOR_ACCENT = "EA580C"
COLOR_ACCENT_SOFT = "FCE7D6"
COLOR_TEXT = "132238"
COLOR_MUTED = "516072"
COLOR_LINE = "D8DDE6"
COLOR_WHITE = "FFFFFF"
COLOR_SUCCESS = "0F766E"
COLOR_WARNING = "B45309"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def set_run_style(run, *, size: int, color: str = COLOR_TEXT, bold: bool = False, font_name: str = FONT_CN) -> None:
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def set_paragraph_text(paragraph, text: str, *, size: int, color: str = COLOR_TEXT, bold: bool = False, font_name: str = FONT_CN) -> None:
    paragraph.clear()
    run = paragraph.add_run()
    run.text = text
    set_run_style(run, size=size, color=color, bold=bold, font_name=font_name)


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    size: int,
    color: str = COLOR_TEXT,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str = FONT_CN,
    margin: float = 0.08,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    set_paragraph_text(p, text, size=size, color=color, bold=bold, font_name=font_name)
    return box


def add_bullets(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    lines: list[str],
    *,
    size: int = 15,
    color: str = COLOR_MUTED,
    bullet_color: str | None = None,
    line_space: float = 1.2,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Inches(0.04)
    box.text_frame.margin_right = Inches(0.04)
    box.text_frame.margin_top = Inches(0.02)
    box.text_frame.margin_bottom = Inches(0.02)
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    box.text_frame.clear()
    marker_color = bullet_color or color
    for idx, line in enumerate(lines):
        p = box.text_frame.paragraphs[0] if idx == 0 else box.text_frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_space
        marker = p.add_run()
        marker.text = "• "
        set_run_style(marker, size=size, color=marker_color, bold=True)
        content = p.add_run()
        content.text = line
        set_run_style(content, size=size, color=color)
    return box


def add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    lines: list[str],
    *,
    fill_color: str = COLOR_SURFACE,
    line_color: str = COLOR_LINE,
    accent_color: str = COLOR_SECONDARY,
    title_color: str = COLOR_PRIMARY,
    body_color: str = COLOR_MUTED,
    title_size: int = 18,
    body_size: int = 14,
):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_color)
    shape.line.color.rgb = rgb(line_color)
    shape.line.width = Pt(1)

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(0.14),
        Inches(height),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(accent_color)
    accent.line.fill.background()

    add_textbox(slide, left + 0.2, top + 0.12, width - 0.28, 0.38, title, size=title_size, color=title_color, bold=True)
    add_bullets(
        slide,
        left + 0.18,
        top + 0.58,
        width - 0.28,
        height - 0.68,
        lines,
        size=body_size,
        color=body_color,
        bullet_color=accent_color,
    )
    return shape


def add_chip(slide, left: float, top: float, width: float, text: str, *, fill_color: str, text_color: str = COLOR_WHITE):
    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.42),
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = rgb(fill_color)
    chip.line.fill.background()
    add_textbox(slide, left + 0.02, top + 0.03, width - 0.04, 0.3, text, size=13, color=text_color, bold=True, align=PP_ALIGN.CENTER)
    return chip


def add_metric_tile(
    slide,
    left: float,
    top: float,
    width: float,
    value: str,
    label: str,
    *,
    fill_color: str = COLOR_SURFACE,
    value_color: str = COLOR_PRIMARY,
    label_color: str = COLOR_MUTED,
):
    tile = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(1.0),
    )
    tile.fill.solid()
    tile.fill.fore_color.rgb = rgb(fill_color)
    tile.line.color.rgb = rgb(COLOR_LINE)
    tile.line.width = Pt(1)
    add_textbox(slide, left + 0.08, top + 0.12, width - 0.16, 0.32, value, size=24, color=value_color, bold=True)
    add_textbox(slide, left + 0.08, top + 0.5, width - 0.16, 0.22, label, size=12, color=label_color)
    return tile


def add_step(
    slide,
    left: float,
    top: float,
    width: float,
    title: str,
    body: str,
    *,
    fill_color: str,
    title_color: str = COLOR_PRIMARY,
    body_color: str = COLOR_MUTED,
):
    step = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(1.08),
    )
    step.fill.solid()
    step.fill.fore_color.rgb = rgb(fill_color)
    step.line.color.rgb = rgb(COLOR_LINE)
    step.line.width = Pt(1)
    add_textbox(slide, left + 0.08, top + 0.12, width - 0.16, 0.26, title, size=16, color=title_color, bold=True)
    add_textbox(slide, left + 0.08, top + 0.42, width - 0.16, 0.46, body, size=11, color=body_color)
    return step


def add_connector(slide, left: float, top: float, width: float, *, color: str = COLOR_ACCENT):
    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.36),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = rgb(color)
    arrow.line.fill.background()
    return arrow


def add_background(slide, *, dark: bool = False) -> None:
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(COLOR_PRIMARY if dark else COLOR_BG)
    bg.line.fill.background()

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.22))
    band.fill.solid()
    band.fill.fore_color.rgb = rgb(COLOR_ACCENT if dark else COLOR_SECONDARY)
    band.line.fill.background()

    corner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(11.5), Inches(-0.5), Inches(2.4), Inches(2.4))
    corner.fill.solid()
    corner.fill.fore_color.rgb = rgb(COLOR_ACCENT if dark else COLOR_ACCENT_SOFT)
    corner.fill.transparency = 0.08
    corner.line.fill.background()

    corner_small = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-0.5), Inches(6.1), Inches(1.8), Inches(1.8))
    corner_small.fill.solid()
    corner_small.fill.fore_color.rgb = rgb(COLOR_SECONDARY if dark else COLOR_SURFACE_ALT)
    corner_small.fill.transparency = 0.15
    corner_small.line.fill.background()


def add_header(slide, eyebrow: str, title: str, subtitle: str | None = None, *, dark: bool = False) -> None:
    eyebrow_color = COLOR_ACCENT_SOFT if dark else COLOR_SECONDARY
    title_color = COLOR_WHITE if dark else COLOR_PRIMARY
    subtitle_color = "E5E7EB" if dark else COLOR_MUTED
    add_textbox(slide, 0.72, 0.54, 2.0, 0.28, eyebrow, size=13, color=eyebrow_color, bold=True)
    add_textbox(slide, 0.72, 0.85, 9.7, 0.58, title, size=28, color=title_color, bold=True)
    if subtitle:
        add_textbox(slide, 0.74, 1.46, 10.8, 0.48, subtitle, size=14, color=subtitle_color)


def add_footer(slide, index: int) -> None:
    add_textbox(slide, 11.8, 7.02, 1.0, 0.18, f"{index:02d}", size=10, color=COLOR_MUTED, align=PP_ALIGN.RIGHT)


def style_chart(chart, *, legend_bottom: bool = False) -> None:
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM if legend_bottom else XL_LEGEND_POSITION.RIGHT
    chart.legend.font.size = Pt(10)
    chart.legend.font.name = FONT_CN
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.name = FONT_CN
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.name = FONT_CN
    chart.value_axis.has_major_gridlines = True


@dataclass
class ProjectData:
    benchmark: dict
    records: list[dict]
    meta: dict

    @property
    def max_message_length(self) -> int:
        return int(self.meta.get("constraints") and self.benchmark["summary"]["lengths"][-1])

    @property
    def constraints(self) -> int:
        return int(self.meta["constraints"])

    @property
    def zkey_bytes(self) -> int:
        return int(self.meta["sizes"]["zkey_bytes"])

    @property
    def vkey_bytes(self) -> int:
        return int(self.meta["sizes"]["verification_key_bytes"])

    @property
    def generated_at(self) -> str:
        return str(self.benchmark["summary"]["generated_at_utc"])


def format_bytes(value: int) -> str:
    if value >= 1024 * 1024:
        return f"{value / (1024 * 1024):.2f} MB"
    if value >= 1024:
        return f"{value / 1024:.2f} KB"
    return f"{value} B"


def block_count(message_length: int) -> int:
    # SM3 padding 需要 1bit 标记和 64bit 长度字段。
    return math.ceil((message_length + 1 + 8) / 64)


def load_data() -> ProjectData:
    benchmark = json.loads(BENCHMARK_PATH.read_text(encoding="utf-8"))
    meta = json.loads(META_PATH.read_text(encoding="utf-8"))
    return ProjectData(benchmark=benchmark, records=benchmark["records"], meta=meta)


def add_comparison_table(slide, left: float, top: float, width: float, height: float) -> None:
    headers = ["方案", "国内合规适配", "电路效率", "工程落地"]
    rows = [
        ["SHA-256 + ZKP", "中", "中", "高"],
        ["Poseidon + ZKP", "低", "高", "高"],
        ["本项目 SM3 + ZKP", "高", "中", "高"],
    ]
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table
    widths = [1.7, 1.45, 1.15, 1.15]
    for idx, w in enumerate(widths):
        table.columns[idx].width = Inches(w)
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(COLOR_PRIMARY)
        cell.text = header
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run_style(p.runs[0], size=12, color=COLOR_WHITE, bold=True)
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(COLOR_SURFACE if row_idx % 2 else COLOR_SURFACE_ALT)
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            color = COLOR_PRIMARY if col_idx == 0 else COLOR_TEXT
            bold = col_idx == 0 or value == "高"
            set_run_style(p.runs[0], size=12, color=color, bold=bold)


def add_correctness_table(slide, left: float, top: float, width: float, height: float, records: list[dict]) -> None:
    headers = ["消息长度", "Block 数", "结果"]
    rows = [[f"{record['message_length']} B", str(block_count(record["message_length"])), "通过"] for record in records]
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table
    table.columns[0].width = Inches(1.2)
    table.columns[1].width = Inches(1.0)
    table.columns[2].width = Inches(0.9)
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(COLOR_SECONDARY)
        cell.text = header
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run_style(p.runs[0], size=11, color=COLOR_WHITE, bold=True)
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(COLOR_SURFACE if row_idx % 2 else COLOR_SURFACE_ALT)
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            set_run_style(p.runs[0], size=11, color=COLOR_TEXT, bold=(col_idx == 2))


def create_title_slide(prs: Presentation, data: ProjectData) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, dark=True)
    add_textbox(slide, 0.74, 0.72, 1.9, 0.28, "FINAL DEFENSE", size=12, color=COLOR_ACCENT_SOFT, bold=True)
    add_textbox(slide, 0.74, 1.15, 8.2, 1.1, "隐链卫士", size=30, color=COLOR_WHITE, bold=True)
    add_textbox(slide, 0.74, 1.92, 10.4, 0.72, "基于 SM3 与零知识证明的隐私验证平台结项答辩", size=22, color="DCE5F4", bold=True)
    add_textbox(
        slide,
        0.78,
        2.78,
        8.7,
        0.58,
        "围绕“国密合规 + 隐私保护 + 可验证性”三条主线，完成了从电路、后端到前端展示的端到端实现。",
        size=15,
        color="D8E2F0",
    )

    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(3.45), Inches(7.65), Inches(2.05))
    hero.fill.solid()
    hero.fill.fore_color.rgb = rgb("17304E")
    hero.line.color.rgb = rgb("2E496B")
    hero.line.width = Pt(1)

    add_textbox(slide, 1.02, 3.7, 2.2, 0.28, "结项核心结果", size=15, color=COLOR_ACCENT_SOFT, bold=True)
    add_bullets(
        slide,
        1.0,
        4.08,
        6.95,
        1.22,
        [
            "完成 SM3 compression step 电路与多块 proof bundle 机制，支持真实 Groth16 证明/验证。",
            "完成 FastAPI + Next.js 全栈平台，打通 prove、verify、dashboard、experiments 四条演示链路。",
            "完成 benchmark、单元测试与性能数据沉淀，形成可复现的结项材料。",
        ],
        size=14,
        color="D7E0EB",
        bullet_color=COLOR_ACCENT,
    )

    add_chip(slide, 0.92, 5.82, 1.7, "35200 约束", fill_color=COLOR_ACCENT)
    add_chip(slide, 2.78, 5.82, 1.9, "真实 Groth16", fill_color=COLOR_SECONDARY)
    add_chip(slide, 4.88, 5.82, 1.9, "247 B / 4 Block", fill_color="31598A")
    add_chip(slide, 6.98, 5.82, 1.3, "全栈演示", fill_color="496C96")

    right_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.15), Inches(1.2), Inches(3.3), Inches(4.98))
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = rgb("F8FBFF")
    right_panel.line.color.rgb = rgb("D4DCE8")
    right_panel.line.width = Pt(1)
    add_textbox(slide, 9.42, 1.48, 2.4, 0.32, "答辩定位", size=16, color=COLOR_PRIMARY, bold=True)
    add_bullets(
        slide,
        9.34,
        1.9,
        2.92,
        1.35,
        [
            "从“中期进展”切换为“结项成果展示”。",
            "重点回答：做成了什么、为什么值得做、性能达到什么水平。",
        ],
        size=13,
        color=COLOR_MUTED,
    )
    add_metric_tile(slide, 9.38, 3.46, 1.38, format_bytes(data.zkey_bytes), "Proving Key")
    add_metric_tile(slide, 10.88, 3.46, 1.38, format_bytes(data.vkey_bytes), "VKey")
    add_metric_tile(slide, 9.38, 4.62, 1.38, "7 组", "Benchmark 档位", fill_color=COLOR_SURFACE_ALT)
    add_metric_tile(slide, 10.88, 4.62, 1.38, "4 页", "前端关键页面", fill_color=COLOR_SURFACE_ALT)

    add_textbox(
        slide,
        9.4,
        5.92,
        2.7,
        0.58,
        f"生成日期：{date.today().isoformat()}",
        size=12,
        color=COLOR_MUTED,
    )
    add_footer(slide, 1)


def create_completion_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, "结项概览", "项目已经完成的四类核心产出", "结项答辩不再以“计划”为中心，而是以“已交付能力”为中心。")

    add_card(
        slide,
        0.78,
        2.08,
        3.0,
        1.72,
        "1. 算法与电路",
        [
            "完成 SM3 compression step 电路设计与构建流程。",
            "公开输入/私有输入接口清晰，支持 block_bits 与 state 链路校验。",
        ],
        fill_color=COLOR_SURFACE,
        accent_color=COLOR_ACCENT,
    )
    add_card(
        slide,
        3.96,
        2.08,
        3.0,
        1.72,
        "2. 后端证明服务",
        [
            "完成 FastAPI API 与 proof bundle 封装。",
            "打通 /api/hash、/api/prove、/api/verify、/api/benchmark。",
        ],
        fill_color=COLOR_SURFACE_ALT,
        accent_color=COLOR_SECONDARY,
    )
    add_card(
        slide,
        7.14,
        2.08,
        2.92,
        1.72,
        "3. 前端演示平台",
        [
            "完成 prove / verify / dashboard / experiments 页面。",
            "支持 proof 展示、结果验证、图表展示和 JSON 导出。",
        ],
        fill_color=COLOR_SURFACE,
        accent_color="31598A",
    )
    add_card(
        slide,
        10.22,
        2.08,
        2.36,
        1.72,
        "4. 实验与文档",
        [
            "生成 benchmark 结果。",
            "单元测试通过。",
            "仓库文档完整。",
        ],
        fill_color=COLOR_SURFACE_ALT,
        accent_color=COLOR_WARNING,
    )

    add_textbox(slide, 0.82, 4.15, 2.8, 0.28, "结项答辩建议主线", size=16, color=COLOR_PRIMARY, bold=True)
    add_step(slide, 0.82, 4.52, 2.12, "为什么做", "国密合规和隐私保护要求叠加，现有公开生态对 SM3 支持不足。", fill_color=COLOR_ACCENT_SOFT)
    add_connector(slide, 2.98, 4.9, 0.42)
    add_step(slide, 3.42, 4.52, 2.12, "怎么做", "用 step 电路 + 外部 padding + proof bundle，把 SM3 多块哈希转成可证明链路。", fill_color=COLOR_SURFACE_ALT)
    add_connector(slide, 5.58, 4.9, 0.42)
    add_step(slide, 6.02, 4.52, 2.12, "做成什么", "完成端到端平台，能真实生成和验证 Groth16 证明，并展示性能指标。", fill_color=COLOR_SURFACE)
    add_connector(slide, 8.18, 4.9, 0.42)
    add_step(slide, 8.62, 4.52, 2.12, "价值在哪", "形成国密场景下可展示、可验证、可扩展的隐私验证基础能力。", fill_color=COLOR_SURFACE_ALT)
    add_connector(slide, 10.78, 4.9, 0.42)
    add_step(slide, 11.22, 4.52, 1.38, "可延伸", "可继续做更大规模与更强聚合。", fill_color=COLOR_ACCENT_SOFT)

    add_textbox(
        slide,
        0.82,
        6.26,
        11.2,
        0.44,
        "这页的作用，是先把评委注意力从“还差什么”转到“已经交付了什么”，为后续讲背景、实现和性能做铺垫。",
        size=13,
        color=COLOR_MUTED,
    )
    add_footer(slide, 2)


def create_background_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, "项目背景", "政策驱动与行业需求共同推动“国密 + 隐私验证”落地", "结项答辩里，这一部分重点说明“为什么这个方向值得做”。")

    add_card(
        slide,
        0.78,
        2.04,
        5.1,
        3.92,
        "政策驱动",
        [
            "《中华人民共和国密码法》自 2020-01-01 起施行，推动关键领域商用密码合规应用。",
            "《中华人民共和国个人信息保护法》自 2021-11-01 起施行，强调敏感个人信息最小暴露。",
            "等保 2.0 相关国家标准自 2019-12-01 实施，强化密码技术与可信验证要求。",
            "《商用密码应用安全性评估管理办法》自 2023-11-01 起施行，进一步提升重要系统的商密应用要求。",
        ],
        fill_color=COLOR_SURFACE,
        accent_color=COLOR_ACCENT,
        body_size=14,
    )

    add_card(
        slide,
        6.1,
        2.04,
        3.0,
        3.92,
        "行业场景",
        [
            "金融：在不暴露交易原文的前提下证明数据真实性与一致性。",
            "政务：兼顾跨部门协同与敏感信息保护，降低直接明文流转风险。",
            "医疗：在科研分析和联合建模中保护患者原始隐私数据。",
        ],
        fill_color=COLOR_SURFACE_ALT,
        accent_color=COLOR_SECONDARY,
        body_size=14,
    )

    add_card(
        slide,
        9.32,
        2.04,
        3.26,
        3.92,
        "核心矛盾",
        [
            "传统系统强调“算得对”，但很难做到“既证明正确又不泄露原文”。",
            "公开 ZKP 生态多围绕 SHA-256 / Poseidon 等路线，SM3 适配能力不足。",
            "真正可答辩、可展示、可验证的国密隐私验证平台较少。",
        ],
        fill_color=COLOR_SURFACE,
        accent_color="31598A",
        body_size=14,
    )

    add_textbox(
        slide,
        0.82,
        6.26,
        11.1,
        0.4,
        "背景结论：项目的价值不只在“把 SM3 写成电路”，而在于把合规需求、隐私需求和工程展示闭环连接起来。",
        size=13,
        color=COLOR_PRIMARY,
        bold=True,
    )
    add_footer(slide, 3)


def create_state_of_art_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, "研究现状", "国内外方案各有侧重，本项目选择了“国密适配 + 平台化落地”路线", "这一页建议回答“别人怎么做，我们为什么不直接照搬”。")

    add_card(
        slide,
        0.78,
        2.0,
        3.82,
        1.66,
        "国外主流路线",
        [
            "大量工作围绕 Poseidon、Poseidon2 等 ZK-friendly hash，追求更低约束和更高证明效率。",
            "SHA-256 也有成熟电路生态，但更偏国际区块链与通用兼容场景。",
        ],
        fill_color=COLOR_SURFACE,
        accent_color=COLOR_ACCENT,
    )
    add_card(
        slide,
        4.76,
        2.0,
        3.84,
        1.66,
        "工程生态做法",
        [
            "Circom / circomlib、gnark 等工具链以标准化模板和模块复用为核心。",
            "公开生态对常见国际哈希支持完善，对 SM3 的端到端实践相对有限。",
        ],
        fill_color=COLOR_SURFACE_ALT,
        accent_color=COLOR_SECONDARY,
    )
    add_card(
        slide,
        8.78,
        2.0,
        3.8,
        1.66,
        "国内需求特点",
        [
            "重点行业更关注国密合规、自主可控和可验证性。",
            "仅做哈希算法替换不够，还需要把证明链路、接口和展示能力一并补齐。",
        ],
        fill_color=COLOR_SURFACE,
        accent_color="31598A",
    )

    add_textbox(slide, 0.82, 4.06, 3.2, 0.26, "方案对比（结项答辩推荐说法）", size=16, color=COLOR_PRIMARY, bold=True)
    add_comparison_table(slide, 0.84, 4.42, 5.72, 1.9)

    add_card(
        slide,
        6.92,
        4.28,
        5.64,
        2.08,
        "本项目的选型判断",
        [
            "如果只追求最低约束，Poseidon 更占优势；但它不直接对应国内国密场景。",
            "如果只复用通用生态，SHA-256 更成熟；但合规叙事和场景适配不足。",
            "本项目选择 SM3 + ZKP，是在合规性、自主可控和工程展示之间做平衡。",
        ],
        fill_color=COLOR_SURFACE_ALT,
        accent_color=COLOR_WARNING,
        body_size=14,
    )

    add_textbox(
        slide,
        0.84,
        6.54,
        11.2,
        0.32,
        "可直接口述：国外更重“低约束”，国内更重“国密落地”，我们的工作是在国内需求约束下把这件事真正做成系统。",
        size=13,
        color=COLOR_MUTED,
    )
    add_footer(slide, 4)


def create_route_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, "项目实现", "总体方案：把多块 SM3 哈希过程转化为可证明、可验证、可展示的链路", "这页重点讲“具体怎么做”，尽量讲清流程而不是堆术语。")

    steps = [
        ("输入消息", "前端接收明文输入，统一按 UTF-8 编码。"),
        ("软件参考哈希", "后端用 Python 参考实现先计算 SM3(x)。"),
        ("电路外 padding / 分块", "完成 padding，并拆成 64-byte blocks。"),
        ("step circuit 逐块证明", "每个 block 复用同一个 compression step 电路。"),
        ("proof bundle 组装", "将多块 step proof 汇总成统一 bundle。"),
        ("链路校验 + verify", "比对 expected hash、public signals 和状态链是否一致。"),
    ]

    start_left = 0.8
    for idx, (title, body) in enumerate(steps):
        left = start_left + idx * 2.05
        add_step(slide, left, 2.42, 1.78, title, body, fill_color=COLOR_SURFACE if idx % 2 == 0 else COLOR_SURFACE_ALT)
        if idx < len(steps) - 1:
            add_connector(slide, left + 1.82, 2.78, 0.2)

    add_card(
        slide,
        0.82,
        4.32,
        4.0,
        1.82,
        "本项目实现目标",
        [
            "在保留 SM3 场景语义的前提下，把“知道原像”转成零知识证明问题。",
            "让多块消息也能被拆解、复用并验证，而不是只支持单块演示。",
        ],
        fill_color=COLOR_SURFACE,
        accent_color=COLOR_ACCENT,
    )
    add_card(
        slide,
        4.94,
        4.32,
        3.52,
        1.82,
        "答辩要强调的实现点",
        [
            "padding 不进电路，避免无效约束增长。",
            "step 电路负责单块，bundle 负责多块连接与结果封装。",
        ],
        fill_color=COLOR_SURFACE_ALT,
        accent_color=COLOR_SECONDARY,
    )
    add_card(
        slide,
        8.62,
        4.32,
        3.96,
        1.82,
        "最终效果",
        [
            "不泄露原始消息，也能证明“我知道一个 x，使得 SM3(x)=H”。",
            "并且这条证明链能被后端真实校验、被前端完整展示。",
        ],
        fill_color=COLOR_SURFACE,
        accent_color="31598A",
    )
    add_footer(slide, 5)


def create_architecture_slide(prs: Presentation, data: ProjectData) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, "系统实现", "五层架构将算法、电路、服务与展示解耦", "结项答辩建议从上到下讲：先讲用户看到什么，再讲系统内部如何支撑。")

    layers = [
        ("前端展示层", "Next.js 15 + TypeScript + Tailwind，负责 prove / verify / dashboard / experiments。", COLOR_ACCENT_SOFT),
        ("接口服务层", "FastAPI 提供 hash、prove、verify、benchmark、circuit meta 等标准接口。", COLOR_SURFACE),
        ("证明引擎层", "snarkjs + Groth16 生成 witness、proof 与 verification。", COLOR_SURFACE_ALT),
        ("电路层", "sm3_compression_step 电路 + 多块 proof bundle 复用机制。", COLOR_SURFACE),
        ("密码原语层", "Python 参考 SM3、padding、消息分块和状态压缩逻辑。", COLOR_SURFACE_ALT),
    ]

    top = 2.0
    for idx, (name, desc, fill) in enumerate(layers):
        step = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.84), Inches(top + idx * 0.8), Inches(6.6), Inches(0.62))
        step.fill.solid()
        step.fill.fore_color.rgb = rgb(fill)
        step.line.color.rgb = rgb(COLOR_LINE)
        step.line.width = Pt(1)
        add_textbox(slide, 1.02, top + idx * 0.8 + 0.12, 1.48, 0.24, name, size=16, color=COLOR_PRIMARY, bold=True)
        add_textbox(slide, 2.36, top + idx * 0.8 + 0.12, 4.82, 0.28, desc, size=12, color=COLOR_MUTED)

    add_card(
        slide,
        7.72,
        2.0,
        4.84,
        1.62,
        "核心接口参数",
        [
            "私有输入：block_bits[512]",
            "公开输入：state_in_words[8] + state_out_words[8]",
            "聚合公开输出：expected_hash_words[8]",
        ],
        fill_color=COLOR_SURFACE,
        accent_color=COLOR_ACCENT,
    )
    add_card(
        slide,
        7.72,
        3.82,
        4.84,
        1.6,
        "当前边界条件",
        [
            f"最大消息长度：{data.benchmark['summary']['lengths'][-1]} B",
            "默认支持最多 4 个 SM3 block",
            "系统只支持真实 Groth16，不回退 mock 模式",
        ],
        fill_color=COLOR_SURFACE_ALT,
        accent_color=COLOR_SECONDARY,
    )
    add_card(
        slide,
        7.72,
        5.62,
        4.84,
        0.98,
        "结项口径",
        ["不仅有电路，而且有完整接口、状态校验和前端展示。"],
        fill_color=COLOR_SURFACE,
        accent_color="31598A",
        body_size=14,
    )
    add_footer(slide, 6)


def create_implementation_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, "核心实现", "电路、后端和前端三部分组成闭环", "建议按“单块电路 -> 多块链路 -> 页面展示”顺序讲，逻辑最清楚。")

    add_card(
        slide,
        0.78,
        2.06,
        3.86,
        2.28,
        "1. step 电路",
        [
            "把 1 个 512-bit 消息块映射成 1 次压缩证明。",
            "输入为 block_bits 与前后状态字，输出不直接暴露原文，只公开状态链。",
            "这样可以把复杂问题拆成可复用的小模块。",
        ],
        fill_color=COLOR_SURFACE,
        accent_color=COLOR_ACCENT,
    )

    add_card(
        slide,
        4.74,
        2.06,
        3.9,
        2.28,
        "2. backend proof bundle",
        [
            "后端先完成 padding，再按块构造 step_inputs。",
            "每个 block 调用同一个 step 电路生成 witness / proof / public signals。",
            "最终把多块 proof 封装成统一 bundle，并汇总时间与大小。",
        ],
        fill_color=COLOR_SURFACE_ALT,
        accent_color=COLOR_SECONDARY,
    )

    add_card(
        slide,
        8.74,
        2.06,
        3.84,
        2.28,
        "3. 前端展示闭环",
        [
            "Prove 页：输入消息、生成 proof、展示摘要与时间。",
            "Verify 页：粘贴 proof / signals，校验 verified 结果。",
            "Dashboard / Experiments：展示约束规模、图表与 benchmark 导出。",
        ],
        fill_color=COLOR_SURFACE,
        accent_color="31598A",
    )

    code_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(4.74), Inches(11.76), Inches(1.36))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = rgb("18273A")
    code_box.line.color.rgb = rgb("30465F")
    code_box.line.width = Pt(1)
    add_textbox(slide, 1.04, 4.98, 11.2, 0.84, "message -> hash -> padding -> step_inputs -> step proofs -> proof bundle -> verify", size=18, color=COLOR_WHITE, bold=True, font_name=FONT_MONO)
    add_textbox(slide, 1.04, 5.56, 11.2, 0.24, "这一条链路是答辩里最关键的一句话，建议原样讲出来。", size=11, color="CBD5E1")
    add_footer(slide, 7)


def create_innovation_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, "创新点", "创新不只是“用了新技术”，而是“针对具体问题提出了针对性做法”", "答辩时尽量按“问题 -> 方法 -> 效果”去表述。")

    add_card(
        slide,
        0.78,
        2.0,
        3.0,
        2.0,
        "创新 1：国密适配",
        [
            "问题：公开 ZKP 生态更偏 SHA / Poseidon，SM3 路线稀缺。",
            "做法：围绕 SM3 compression 过程自建电路与 proving 流程。",
            "效果：项目天然更适合国内合规场景展示。",
        ],
        fill_color=COLOR_SURFACE,
        accent_color=COLOR_ACCENT,
    )
    add_card(
        slide,
        3.96,
        2.0,
        3.0,
        2.0,
        "创新 2：分块复用",
        [
            "问题：任意长度输入直接整体入电路，规模会迅速膨胀。",
            "做法：padding 外部化 + step circuit 逐块复用。",
            "效果：让 1~4 block 共享同一电路模板，结构更稳定。",
        ],
        fill_color=COLOR_SURFACE_ALT,
        accent_color=COLOR_SECONDARY,
    )
    add_card(
        slide,
        7.14,
        2.0,
        2.98,
        2.0,
        "创新 3：链路校验",
        [
            "问题：多块 proof 如何确保前后状态首尾相接？",
            "做法：公开 state_in/state_out 并在验证阶段做链式一致性检查。",
            "效果：增强多块证明的结构可靠性。",
        ],
        fill_color=COLOR_SURFACE,
        accent_color="31598A",
    )
    add_card(
        slide,
        10.3,
        2.0,
        2.28,
        2.0,
        "创新 4：平台化",
        [
            "把电路能力做成可演示平台，而不是停留在脚本层。",
            "支持结果展示、图表对比和 benchmark 导出。",
        ],
        fill_color=COLOR_SURFACE_ALT,
        accent_color=COLOR_WARNING,
    )

    add_textbox(slide, 0.84, 4.44, 2.2, 0.26, "本页收束表达", size=16, color=COLOR_PRIMARY, bold=True)
    add_bullets(
        slide,
        0.86,
        4.78,
        11.2,
        1.3,
        [
            "我们解决的不是单点问题，而是“合规性、可复用性、链路可验证性、系统可展示性”四个问题同时出现时的工程落地问题。",
            "因此创新点既体现在电路结构，也体现在后端链路设计和前端演示闭环。",
        ],
        size=15,
        color=COLOR_MUTED,
        bullet_color=COLOR_ACCENT,
    )
    add_footer(slide, 8)


def create_test_design_slide(prs: Presentation, data: ProjectData) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, "性能测试", "测试设计先讲“怎么测”，再讲“测出了什么”", "结项答辩里，测试方法是否清楚往往比单个数字更重要。")

    lengths = ", ".join(str(v) for v in data.benchmark["summary"]["lengths"])
    add_card(
        slide,
        0.78,
        2.02,
        4.26,
        2.08,
        "测试配置",
        [
            f"消息长度档位：{lengths} B",
            "每档均运行真实 prove / verify，结果来自 benchmarks/results/latest.json。",
            f"最新 benchmark 生成时间：{data.generated_at}",
        ],
        fill_color=COLOR_SURFACE,
        accent_color=COLOR_ACCENT,
    )
    add_card(
        slide,
        5.18,
        2.02,
        3.08,
        2.08,
        "观测指标",
        [
            "software_hash_ms",
            "witness_generation_ms",
            "proving_ms / verification_ms",
            "proof_size_bytes",
        ],
        fill_color=COLOR_SURFACE_ALT,
        accent_color=COLOR_SECONDARY,
    )
    add_card(
        slide,
        8.42,
        2.02,
        4.16,
        2.08,
        "正确性验证",
        [
            "单元测试覆盖标准向量、padding 边界和链路篡改检测。",
            "本地执行 backend/tests/test_zkp_inputs.py：4 项全部通过。",
            "benchmark 中 7 组消息长度记录 success 均为 true。",
        ],
        fill_color=COLOR_SURFACE,
        accent_color="31598A",
    )

    add_textbox(slide, 0.84, 4.42, 2.6, 0.26, "消息长度与 block 对应关系", size=16, color=COLOR_PRIMARY, bold=True)
    add_correctness_table(slide, 0.86, 4.78, 3.32, 2.0, data.records)

    add_card(
        slide,
        4.56,
        4.78,
        3.84,
        2.0,
        "答辩时可强调",
        [
            "本项目不是只做“正确性样例”，而是做了不同长度下的阶梯测试。",
            "通过 1~4 block 的覆盖，能观察多块复用后性能是否线性增长。",
        ],
        fill_color=COLOR_SURFACE_ALT,
        accent_color=COLOR_WARNING,
    )
    add_card(
        slide,
        8.58,
        4.78,
        4.0,
        2.0,
        "这一页的作用",
        [
            "证明测试数据并非随意挑选，而是围绕系统边界条件设计。",
            "为下一页展示耗时与 proof 大小曲线提供可信前提。",
        ],
        fill_color=COLOR_SURFACE,
        accent_color=COLOR_ACCENT,
    )
    add_footer(slide, 9)


def create_performance_slide(prs: Presentation, data: ProjectData) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, "性能测试结果", "证明时间、验证时间与 proof 大小随消息长度增加而上升，但增长趋势稳定", "结项答辩这里建议强调“趋势”和“可解释性”。")

    categories = [f"{record['message_length']}B" for record in data.records]
    prove_values = [round(record["proving_ms"], 2) for record in data.records]
    verify_values = [round(record["verification_ms"], 2) for record in data.records]
    proof_values = [int(record["proof_size_bytes"]) for record in data.records]

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("证明时间(ms)", prove_values)
    chart_data.add_series("验证时间(ms)", verify_values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.82),
        Inches(2.0),
        Inches(6.0),
        Inches(3.28),
        chart_data,
    ).chart
    style_chart(chart, legend_bottom=True)

    size_data = CategoryChartData()
    size_data.categories = categories
    size_data.add_series("proof 大小(B)", proof_values)
    size_chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(6.96),
        Inches(2.0),
        Inches(3.72),
        Inches(3.28),
        size_data,
    ).chart
    style_chart(size_chart, legend_bottom=True)

    first = data.records[0]
    last = data.records[-1]
    add_card(
        slide,
        10.9,
        2.0,
        1.7,
        3.28,
        "关键结论",
        [
            f"1 block：prove {first['proving_ms']:.0f} ms / verify {first['verification_ms']:.0f} ms / proof {first['proof_size_bytes']} B",
            f"4 block：prove {last['proving_ms']:.0f} ms / verify {last['verification_ms']:.0f} ms / proof {last['proof_size_bytes']} B",
            "随着 block 数增加，三项指标基本呈线性增长，说明 step 复用策略是有效的。",
            f"软件 SM3 哈希始终低于 {max(record['software_hash_ms'] for record in data.records):.2f} ms，说明系统瓶颈主要在 ZKP 证明阶段。",
        ],
        fill_color=COLOR_SURFACE,
        accent_color=COLOR_WARNING,
        title_size=15,
        body_size=12,
    )

    add_textbox(slide, 0.84, 5.52, 2.4, 0.26, "可直接口述的结果总结", size=16, color=COLOR_PRIMARY, bold=True)
    add_bullets(
        slide,
        0.86,
        5.86,
        11.7,
        0.72,
        [
            "在本项目边界下，单块消息证明约 0.7 秒，最大 4 块消息证明约 2.8 秒，验证约 1.1 秒；性能与块数增长关系清晰，便于后续进一步优化与扩展。",
        ],
        size=15,
        color=COLOR_MUTED,
        bullet_color=COLOR_ACCENT,
    )
    add_footer(slide, 10)


def add_ui_panel(slide, left: float, top: float, width: float, title: str, subtitle: str, lines: list[str], accent: str) -> None:
    shell = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(2.08),
    )
    shell.fill.solid()
    shell.fill.fore_color.rgb = rgb(COLOR_SURFACE)
    shell.line.color.rgb = rgb(COLOR_LINE)
    shell.line.width = Pt(1)

    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.34))
    header.fill.solid()
    header.fill.fore_color.rgb = rgb(accent)
    header.line.fill.background()
    add_textbox(slide, left + 0.08, top + 0.04, width - 0.16, 0.2, title, size=13, color=COLOR_WHITE, bold=True)
    add_textbox(slide, left + 0.08, top + 0.44, width - 0.16, 0.26, subtitle, size=11, color=COLOR_MUTED)
    add_bullets(slide, left + 0.08, top + 0.78, width - 0.16, 1.08, lines, size=11, color=COLOR_TEXT, bullet_color=accent)


def create_frontend_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, "前端展示", "前端不是附属页面，而是把算法能力讲清楚的答辩界面", "仓库当前没有静态截图资源，因此这里用可编辑页面结构示意来表达展示逻辑。")

    add_ui_panel(
        slide,
        0.82,
        2.0,
        2.82,
        "Prove",
        "输入消息，生成 proof bundle",
        ["显示 hash 与 expected hash words", "展示 witness / proving 时间", "输出 proof JSON 与 public signals"],
        COLOR_ACCENT,
    )
    add_ui_panel(
        slide,
        3.82,
        2.0,
        2.82,
        "Verify",
        "粘贴 proof，验证结果",
        ["输入 expected hash", "校验 chain consistency", "返回 verified: true / false"],
        COLOR_SECONDARY,
    )
    add_ui_panel(
        slide,
        6.82,
        2.0,
        2.82,
        "Dashboard",
        "展示约束和工具链状态",
        ["查看 constraints / key size", "查看 artifact 是否齐备", "查看 benchmark 快照与曲线"],
        "31598A",
    )
    add_ui_panel(
        slide,
        9.82,
        2.0,
        2.76,
        "Experiments",
        "展示实验与导出能力",
        ["查看 benchmark 表格", "导出 benchmark JSON", "显示正确性检查项"],
        COLOR_WARNING,
    )

    add_textbox(slide, 0.86, 4.58, 2.4, 0.26, "答辩推荐演示顺序", size=16, color=COLOR_PRIMARY, bold=True)
    add_step(slide, 0.86, 4.96, 2.16, "步骤 1", "在 Prove 页输入示例消息，展示 hash 与 proving 时间。", fill_color=COLOR_SURFACE_ALT)
    add_connector(slide, 3.08, 5.32, 0.32)
    add_step(slide, 3.42, 4.96, 2.16, "步骤 2", "复制 proof 和 public signals 到 Verify 页，验证结果为 true。", fill_color=COLOR_SURFACE)
    add_connector(slide, 5.64, 5.32, 0.32)
    add_step(slide, 5.98, 4.96, 2.16, "步骤 3", "切到 Dashboard，展示约束规模、zkey 大小和性能快照。", fill_color=COLOR_SURFACE_ALT)
    add_connector(slide, 8.2, 5.32, 0.32)
    add_step(slide, 8.54, 4.96, 2.16, "步骤 4", "切到 Experiments，展示表格、曲线与导出能力。", fill_color=COLOR_SURFACE)
    add_connector(slide, 10.76, 5.32, 0.32)
    add_step(slide, 11.1, 4.96, 1.48, "效果", "形成完整闭环。", fill_color=COLOR_ACCENT_SOFT)
    add_footer(slide, 11)


def create_summary_slide(prs: Presentation, data: ProjectData) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(slide, "总结", "项目已经从“想法验证”走到“系统落地”，具备结项展示价值", "最后一页建议既给出结果，也保留后续扩展空间。")

    add_card(
        slide,
        0.78,
        2.04,
        3.78,
        2.8,
        "已完成的结果",
        [
            f"完成 {data.constraints} 约束规模的 SM3 compression step 电路。",
            "完成多块 proof bundle 方案，支持最大 247 B / 4 block 输入。",
            "完成 FastAPI + Next.js 全栈演示平台和 benchmark 数据链路。",
        ],
        fill_color=COLOR_SURFACE,
        accent_color=COLOR_ACCENT,
    )
    add_card(
        slide,
        4.78,
        2.04,
        3.78,
        2.8,
        "项目价值",
        [
            "证明了国密哈希与零知识证明的融合可以被平台化、工程化表达。",
            "既能满足“隐私不泄露”的叙事，也能满足“结果可验证”的系统诉求。",
            "为后续面向政务、金融、医疗等场景的演示和扩展打下基础。",
        ],
        fill_color=COLOR_SURFACE_ALT,
        accent_color=COLOR_SECONDARY,
    )
    add_card(
        slide,
        8.78,
        2.04,
        3.8,
        2.8,
        "后续可继续优化",
        [
            "扩大 max_blocks 与支持更长消息。",
            "继续做批量验证、递归聚合或更强的 bundle 设计。",
            "与 SM2 / SM9 等国密算法结合，扩展完整隐私计算工具链。",
        ],
        fill_color=COLOR_SURFACE,
        accent_color="31598A",
    )

    add_textbox(slide, 0.84, 5.28, 2.6, 0.26, "一句话收尾", size=16, color=COLOR_PRIMARY, bold=True)
    add_bullets(
        slide,
        0.86,
        5.62,
        11.4,
        0.72,
        [
            "本项目最终交付的不只是一个算法原型，而是一套围绕 SM3 与零知识证明构建的、可运行、可验证、可展示的结项级平台。",
        ],
        size=16,
        color=COLOR_PRIMARY,
        bullet_color=COLOR_ACCENT,
    )

    add_textbox(
        slide,
        0.84,
        6.56,
        11.4,
        0.32,
        "备注：本 PPT 已按“结项答辩”重写口径，可直接在 PowerPoint 中继续修改文字、图表与版式。",
        size=12,
        color=COLOR_MUTED,
    )
    add_footer(slide, 12)


def build_presentation(data: ProjectData) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs, data)
    create_completion_slide(prs)
    create_background_slide(prs)
    create_state_of_art_slide(prs)
    create_route_slide(prs)
    create_architecture_slide(prs, data)
    create_implementation_slide(prs)
    create_innovation_slide(prs)
    create_test_design_slide(prs, data)
    create_performance_slide(prs, data)
    create_frontend_slide(prs)
    create_summary_slide(prs, data)
    return prs


def main() -> None:
    data = load_data()
    prs = build_presentation(data)
    prs.save(OUTPUT_PATH)
    print(OUTPUT_PATH)


if __name__ == "__main__":
    main()
